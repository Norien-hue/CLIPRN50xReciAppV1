from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

DARK_GREEN = RGBColor(0x2E, 0x7D, 0x32)
MEDIUM_GREEN = RGBColor(0x38, 0x8E, 0x3C)
LIGHT_GREEN = RGBColor(0x81, 0xC7, 0x84)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xC6, 0x28, 0x28)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
LIGHTER_GREEN = RGBColor(0xF1, 0xF8, 0xE9)
RED_BG = RGBColor(0xFF, 0xEB, 0xEE)

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)


def add_rect(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=16, bold=False,
                color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_center_text(slide, text, top, font_size, bold=False, color=WHITE):
    return add_textbox(slide, Cm(2), top, Cm(29.867), Cm(2), text,
                       font_size, bold, color, PP_ALIGN.CENTER)


def add_accent_bar(slide):
    add_rect(slide, Cm(0), Cm(0), Cm(0.3), Cm(19.05), DARK_GREEN)


def add_title(slide, text):
    return add_textbox(slide, Cm(1.5), Cm(0.6), Cm(31), Cm(1.5), text,
                       28, True, DARK_GREEN)


def add_bullets(slide, left, top, width, height, items, font_size=16):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
    return txBox


def add_green_dots(slide, items, start_top, left=Cm(1.5), gap=Cm(1.8), size=Cm(0.3)):
    for i in range(len(items)):
        y = start_top + gap * i + Cm(0.15)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y, size, size)
        dot.fill.solid()
        dot.fill.fore_color.rgb = LIGHT_GREEN
        dot.line.fill.background()


def format_table_cell(cell, text, font_size=14, bold=False, color=DARK_GRAY,
                      align=PP_ALIGN.CENTER, fill=None):
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.alignment = align
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def make_header_row(table, headers, fill=DARK_GREEN):
    for i, h in enumerate(headers):
        format_table_cell(table.cell(0, i), h, 16, True, WHITE, PP_ALIGN.CENTER, fill)


# ============================================================
# SLIDE 1 — Portada
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s1, Cm(0), Cm(0), Cm(33.867), Cm(19.05), DARK_GREEN)
add_center_text(s1, "Sistema de Reciclaje Inteligente", Cm(6), 36, True)
add_center_text(s1, "Identificaci\u00f3n de productos mediante CLIP ViT-B/32", Cm(8.5), 20)
add_center_text(s1, "Universidad de Alicante \u2014 Interfaces de Usuario \u2014 Curso 2025\u20132026", Cm(14), 14)

# ============================================================
# SLIDE 2 — Arquitectura del Sistema
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(s2)
add_title(s2, "Arquitectura del Sistema")

cols = [
    ("Cliente JavaFX", [
        "Aplicaci\u00f3n de escritorio MVC",
        "Vistas FXML + Controllers",
        "C\u00e1mara OpenCV integrada",
        "Interfaz verde corporativa",
        "Validaci\u00f3n y registro local"
    ]),
    ("API REST Spring Boot", [
        "Servidor EC2 (puerto :3000)",
        "Autenticaci\u00f3n JWT",
        "Endpoints RESTful",
        "Conexi\u00f3n con base de datos",
        "Validaci\u00f3n de usuarios por TAP"
    ]),
    ("Servicio CLIP", [
        "Contenedor Docker en EC2",
        "Modelo ViT-B/32 de OpenAI",
        "Embeddings de 512 dimensiones",
        "Endpoint en puerto :8080",
        "Comunicaci\u00f3n s\u00edncrona con API"
    ])
]

col_w = Cm(9.5)
col_gap = Cm(0.8)
col_start = Cm(1.8)
col_top = Cm(2.5)
col_h = Cm(14.5)

for i, (header, items) in enumerate(cols):
    x = col_start + (col_w + col_gap) * i
    add_rect(s2, x, col_top, col_w, col_h, LIGHTER_GREEN)
    add_textbox(s2, x, col_top + Cm(0.2), col_w, Cm(1.2), header,
                18, True, MEDIUM_GREEN, PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        y = col_top + Cm(2.0) + Cm(2.0) * j
        add_textbox(s2, x + Cm(0.5), y, col_w - Cm(1), Cm(1.5),
                    f"\u2022 {item}", 13, False, DARK_GRAY)

for i in range(2):
    ax = col_start + (col_w + col_gap) * i + col_w + Cm(0.15)
    add_textbox(s2, ax, Cm(9), Cm(1), Cm(1), "\u2192", 24, False, MEDIUM_GREEN, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3 — Cámara Asíncrona
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(s3)
add_title(s3, "C\u00e1mara As\u00edncrona y ProgressBar")

bullets3 = [
    "OpenCV FrameGrabber para captura de video en tiempo real",
    "Square crop: 640\u00d7480 \u2192 480\u00d7480 (recorte cuadrado centrado)",
    "Inicializaci\u00f3n en hilo secundario (background thread)",
    "ProgressBar indeterminada durante el arranque de la c\u00e1mara",
    "AtomicBoolean guard para evitar doble inicializaci\u00f3n",
    "Liberaci\u00f3n correcta de recursos al cerrar la aplicaci\u00f3n"
]
add_bullets(s3, Cm(2.2), Cm(3), Cm(30), Cm(14), bullets3, 17)
add_green_dots(s3, bullets3, Cm(3))

# ============================================================
# SLIDE 4 — Eliminación de Fondo (Flood Fill)
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(s4)
add_title(s4, "Eliminaci\u00f3n de Fondo: Flood Fill en Java Puro")

# Problem
add_textbox(s4, Cm(1.8), Cm(2.5), Cm(14), Cm(1.2), "Problema",
            20, True, RED)
add_rect(s4, Cm(1.8), Cm(3.5), Cm(14), Cm(5.5), RED_BG)
pb = [
    "GrabCut de OpenCV causa crashes JVM",
    "C\u00f3digo nativo inestable en diferentes plataformas",
    "Dependencia de bibliotecas nativas problem\u00e1tica",
    "Errores dif\u00edciles de depurar y reproducir"
]
add_bullets(s4, Cm(2.3), Cm(3.7), Cm(13), Cm(5), pb, 14)

# Solution
add_textbox(s4, Cm(17.5), Cm(2.5), Cm(14), Cm(1.2), "Soluci\u00f3n",
            20, True, MEDIUM_GREEN)
add_rect(s4, Cm(17.5), Cm(3.5), Cm(14), Cm(5.5), LIGHT_GREEN_BG)
sb = [
    "Flood Fill en Java puro (sin c\u00f3digo nativo)",
    "Gradiente estilo Sobel para detecci\u00f3n de bordes",
    "Flood fill desde esquinas con umbral configurable",
    "Median blur 3\u00d73 en m\u00e1scara resultante",
    "Slider de sensibilidad (5\u2013100) para ajuste fino",
    "Sin c\u00f3digo nativo \u2192 sin crashes JVM"
]
add_bullets(s4, Cm(18), Cm(3.7), Cm(13), Cm(5), sb, 14)

# Comparison table
t4 = s4.shapes.add_table(3, 3, Cm(6), Cm(10.5), Cm(20), Cm(5)).table
t4.columns[0].width = Cm(8)
t4.columns[1].width = Cm(6)
t4.columns[2].width = Cm(6)

make_header_row(t4, ["M\u00e9todo", "Velocidad", "Estabilidad"])

data4 = [
    ["GrabCut (OpenCV)", "R\u00e1pido", "Crash JVM"],
    ["Flood Fill (Java)", "Medio", "100% estable"]
]
fills4 = [RED_BG, LIGHT_GREEN_BG]
colors4 = [DARK_GRAY, DARK_GRAY]
special4 = [True, True]  # col 2 special formatting

for ri, row in enumerate(data4):
    for ci, val in enumerate(row):
        bold = (ci == 2)
        color = RED if (ri == 0 and ci == 2) else (MEDIUM_GREEN if (ri == 1 and ci == 2) else DARK_GRAY)
        format_table_cell(t4.cell(ri + 1, ci), val, 14, bold, color, PP_ALIGN.CENTER, fills4[ri])

# ============================================================
# SLIDE 5 — Tabla y Umbral Rojo
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(s5)
add_title(s5, "Tabla de Resultados con Umbral de Confianza")

bullets5 = [
    "TableView con productos ordenados por score descendente",
    "Umbral de confianza: 47% (RED_THRESHOLD = 0.47)",
    "Productos con score < 47% se muestran en rojo + negrita",
    "Selecci\u00f3n bloqueada (selection listener previene selecci\u00f3n)",
    "Campo de b\u00fasqueda para filtrar por nombre de producto"
]
add_bullets(s5, Cm(2), Cm(3), Cm(30), Cm(6.5), bullets5, 17)

# Visual table
t5 = s5.shapes.add_table(6, 4, Cm(4), Cm(10.5), Cm(25), Cm(6)).table
t5.columns[0].width = Cm(5)
t5.columns[1].width = Cm(9)
t5.columns[2].width = Cm(5)
t5.columns[3].width = Cm(6)

make_header_row(t5, ["C\u00f3digo", "Producto", "Score", "Estado"])

rows5 = [
    ["B001", "Botella PET", "92%", "OK"],
    ["L002", "Lat\u00f3n de refresco", "78%", "OK"],
    ["C003", "Cart\u00f3n leche", "65%", "OK"],
    ["P004", "Pl\u00e1stico film", "42%", "BAJO"],
    ["V005", "Vidrio roto", "31%", "BAJO"]
]

for ri, row in enumerate(rows5):
    score = int(row[2].rstrip('%'))
    is_low = score < 47
    for ci, val in enumerate(row):
        color = RED if is_low else DARK_GRAY
        bold = is_low
        fill = RED_BG if is_low else WHITE
        format_table_cell(t5.cell(ri + 1, ci), val, 14, bold, color, PP_ALIGN.CENTER, fill)

# ============================================================
# SLIDE 6 — CLIP ViT-B/32 y Matching Híbrido
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(s6)
add_title(s6, "CLIP ViT-B/32 y Coincidencia H\u00edbrida")

bullets6 = [
    "Modelo: ViT-B/32 (Vision Transformer), embeddings de 512 dimensiones",
    "8 embeddings de productos precomputados para matching r\u00e1pido",
    "Matching h\u00edbrido: imagen + texto (BLEND_WEIGHT = 0.5)",
    "Inferencia ~500ms en instancia t3.small de AWS",
    "Top-5 resultados devueltos ordenados por similitud",
    "Coseno similitud entre embedding de consulta y precomputados"
]
add_bullets(s6, Cm(2.2), Cm(3), Cm(30), Cm(14), bullets6, 17)

# Highlight first bullet with light green bg box
add_rect(s6, Cm(1.8), Cm(2.9), Cm(30.5), Cm(1.6), LIGHT_GREEN_BG)
# Re-add the first bullet text on top
add_textbox(s6, Cm(2.2), Cm(3), Cm(30), Cm(1.5), bullets6[0], 17, False, DARK_GRAY)

add_green_dots(s6, bullets6, Cm(3))

# ============================================================
# SLIDE 7 — Modal Único
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(s7)
add_title(s7, "Modal \u00danico: Confirmaci\u00f3n y Selecci\u00f3n")

flow = [
    ("1. Capturar", "Tomar foto con la c\u00e1mara OpenCV"),
    ("2. Confirmar", "Vista previa en modal con opciones"),
    ("3. Eliminar fondo", "Slider de sensibilidad (opcional)"),
    ("4. Vista IA", "Previsualizaci\u00f3n 224\u00d7224"),
    ("5. Enviar a API", "Imagen \u2192 API REST \u2192 CLIP \u2192 resultados"),
    ("6. Seleccionar", "TAP + tabla de productos + Accept"),
    ("7. Validar", "Validar TAP y registrar reciclaje")
]

bw = Cm(7)
bh = Cm(3.5)
sx = Cm(1.8)
sy = Cm(2.5)
gx = Cm(1.5)
gy = Cm(1.5)

for idx, (step, desc) in enumerate(flow):
    if idx < 4:
        x = sx + (bw + gx) * idx
        y = sy
    else:
        x = sx + (bw + gx) * (idx - 4)
        y = sy + bh + gy

    add_rect(s7, x, y, bw, bh, LIGHT_GREEN_BG)
    add_textbox(s7, x + Cm(0.3), y + Cm(0.2), bw - Cm(0.6), Cm(1),
                step, 14, True, DARK_GREEN)
    add_textbox(s7, x + Cm(0.3), y + Cm(1.2), bw - Cm(0.6), Cm(2),
                desc, 12, False, DARK_GRAY)

# Arrows upper row
for i in range(3):
    ax = sx + (bw + gx) * i + bw + Cm(0.1)
    add_textbox(s7, ax, sy + Cm(1.2), Cm(1), Cm(1), "\u2192", 18, False, MEDIUM_GREEN, PP_ALIGN.CENTER)

# Down arrow from row 1 to row 2
add_textbox(s7, sx + (bw + gx) * 1.5, sy + bh + Cm(0.05), Cm(1), Cm(1),
            "\u2193", 18, False, MEDIUM_GREEN, PP_ALIGN.CENTER)

# Arrows lower row
for i in range(2):
    ax = sx + (bw + gx) * i + bw + Cm(0.1)
    add_textbox(s7, ax, sy + bh + gy + Cm(1.2), Cm(1), Cm(1),
                "\u2192", 18, False, MEDIUM_GREEN, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8 — Problemas Resueltos
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(s8)
add_title(s8, "Problemas y Soluciones")

t8 = s8.shapes.add_table(6, 2, Cm(3), Cm(3), Cm(27), Cm(13)).table
t8.columns[0].width = Cm(11)
t8.columns[1].width = Cm(16)

make_header_row(t8, ["Problema", "Soluci\u00f3n"])

problems = [
    ["GrabCut crashes JVM", "Flood Fill en Java puro (sin c\u00f3digo nativo)"],
    ["UI freeze al iniciar c\u00e1mara", "Background thread + ProgressBar indeterminada"],
    ["Doble clic en bot\u00f3n de c\u00e1mara", "AtomicBoolean guard para evitar doble inicio"],
    ["Selecci\u00f3n de productos con score bajo", "Umbral rojo (47%) + deshabilitado"],
    ["Disco EC2 lleno (/var)", "Docker system prune automatizado"]
]

for ri, row in enumerate(problems):
    fill = WHITE if ri % 2 == 0 else LIGHTER_GREEN
    for ci, val in enumerate(row):
        bold = (ci == 0)
        format_table_cell(t8.cell(ri + 1, ci), val, 14, bold, DARK_GRAY, PP_ALIGN.LEFT, fill)

# ============================================================
# SLIDE 9 — Conclusiones
# ============================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(s9)
add_title(s9, "Conclusiones")

bullets9 = [
    "CLIP ViT-B/32 identifica productos de reciclaje eficazmente",
    "Java puro evita problemas de c\u00f3digo nativo (sin crashes JVM)",
    "Arquitectura escalable basada en API REST + Docker",
    "Interfaz intuitiva y responsiva con feedback visual",
    "Trabajo futuro: m\u00e1s productos, app m\u00f3vil, mejorar bg removal"
]
add_bullets(s9, Cm(2.2), Cm(3), Cm(30), Cm(14), bullets9, 18)

for i in range(len(bullets9)):
    y = Cm(3) + Cm(2) * i + Cm(0.15)
    dot = s9.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1.5), y, Cm(0.4), Cm(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = MEDIUM_GREEN
    dot.line.fill.background()

# ============================================================
# SLIDE 10 — Gracias
# ============================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s10, Cm(0), Cm(0), Cm(33.867), Cm(19.05), DARK_GREEN)
add_center_text(s10, "Gracias", Cm(6), 44, True)
add_center_text(s10, "\u00bfPreguntas?", Cm(9), 24)

# Save
output_path = r'C:\Users\FA506NC\Desktop\Subida de nota\Interfaces_PI\Presentacion.pptx'
prs.save(output_path)
print("PPTX generated successfully")
